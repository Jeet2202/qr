from pptx import Presentation

def extract_ppt_content(file_path):

    prs = Presentation(file_path)

    slides = []

    for index, slide in enumerate(prs.slides):

        title = ""
        content = []
        bullet_points = []
        image_count = 0

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if shape == slide.shapes.title:
                    title = text
                else:
                    content.append(text)

                    if "\n" in text:
                        bullet_points.extend(text.split("\n"))

            if shape.shape_type == 13:
                image_count += 1

        slides.append({
            "slide_number": index + 1,
            "title": title,
            "content": " ".join(content),
            "bullet_points": bullet_points,
            "image_count": image_count
        })

    return slides